from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


def create_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    colors = {
        "primary": RGBColor(255, 141, 40),
        "secondary": RGBColor(255, 184, 77),
        "cream": RGBColor(255, 244, 230),
        "deep_green": RGBColor(47, 72, 88),
        "accent_green": RGBColor(76, 175, 80),
        "accent_dark_green": RGBColor(46, 125, 50),
        "text_dark": RGBColor(74, 74, 74),
        "white": RGBColor(255, 255, 255),
    }

    def add_background(slide, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def apply_gradient(fill, angle, stops):
        fill.gradient()
        fill.gradient_angle = angle
        gradient_stops = fill.gradient_stops
        available = len(gradient_stops)
        if not available:
            return
        effective_stops = stops
        if len(stops) > available:
            effective_stops = [stops[0]] + [stops[-1]] if available == 2 else stops[:available]
        for idx, (position, color_key) in enumerate(effective_stops[:available]):
            gradient_stops[idx].position = position
            gradient_stops[idx].color.rgb = colors[color_key]

    def add_title(slide, text, top=Inches(0.7)):
        title_box = slide.shapes.add_textbox(Inches(0.7), top, prs.slide_width - Inches(1.4), Inches(1.5))
        tf = title_box.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.font.name = "Poppins"
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = colors["deep_green"]

    def add_list(slide, items, top=Inches(2), left=Inches(0.9)):
        box = slide.shapes.add_textbox(left, top, prs.slide_width - Inches(1.8), Inches(4.5))
        tf = box.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(items):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = item
            p.font.name = "Nunito"
            p.font.size = Pt(20)
            p.font.color.rgb = colors["text_dark"]
            p.level = 0
            p.space_after = Pt(6)

    # Slide 1 - Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    apply_gradient(
        background.fill,
        angle=135,
        stops=[(0.0, "cream"), (0.5, "secondary"), (1.0, "primary")],
    )
    background.line.fill.background()

    cover_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.3), prs.slide_width - Inches(2), Inches(4.5)
    )
    cover_box.fill.solid()
    cover_box.fill.fore_color.rgb = colors["white"]
    cover_box.line.color.rgb = colors["secondary"]
    cover_box.line.width = Pt(3)

    text_box = slide.shapes.add_textbox(Inches(1.4), Inches(1.6), prs.slide_width - Inches(2.8), Inches(4))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Asli Bogor"
    p.font.name = "Poppins"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = colors["primary"]

    subtitle = tf.add_paragraph()
    subtitle.text = "Platform digital untuk menemukan, mendukung, dan mempercepat UMKM otentik Bogor"
    subtitle.font.name = "Nunito"
    subtitle.font.size = Pt(24)
    subtitle.font.color.rgb = colors["deep_green"]
    subtitle.space_before = Pt(10)

    cta = tf.add_paragraph()
    cta.text = "Pitch deck • Q4 2025"
    cta.font.name = "Nunito"
    cta.font.size = Pt(18)
    cta.font.color.rgb = colors["text_dark"]
    cta.space_before = Pt(20)

    # Slide 2 - Brand snapshot
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, colors["cream"])
    add_title(slide, "Brand Snapshot & Audiens")
    add_list(
        slide,
        [
            "Misi: membantu UMKM Bogor naik kelas lewat ekosistem digital yang humanis.",
            "Audiens utama: pelaku UMKM kuliner, craft, fashion, dan warga Bogor pecinta produk lokal.",
            "Tone desain: hangat, natural, penuh energi oranye-hijau identitas Asli Bogor.",
        ],
    )

    # Slide 3 - Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, colors["white"])
    add_title(slide, "Tantangan UMKM Bogor Saat Ini")
    add_list(
        slide,
        [
            "Visibilitas digital rendah di tengah kompetisi marketplace nasional.",
            "Operasional pesanan, pembayaran, dan logistik masih manual.",
            "Kurangnya data pelanggan untuk retensi dan promosi personal.",
            "Driver lokal belum terhubung ke kebutuhan distribusi UMKM.",
        ],
    )

    # Slide 4 - Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, colors["cream"])
    add_title(slide, "Solusi Platform Asli Bogor")
    add_list(
        slide,
        [
            "Direktori UMKM terkurasi dengan visual hangat sesuai kategori.",
            "Autentikasi multi-aktor dan onboarding cepat (UMKM, user, driver).",
            "Dashboard analytics, notifikasi, dan konten edukasi untuk UMKM.",
            "Integrasi pesan-antar driver lokal menjaga pengalaman end-to-end.",
        ],
    )

    # Slide 5 - Features
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, colors["white"])
    add_title(slide, "Fitur Unggulan Website")
    feature_cards = [
        ("Produk Lokal Terkurasi", "Card kategori makanan, minuman, fashion, jasa dengan warna gradien."),
        ("Pesan Antar Cepat", "CTA jelas dan dukungan driver lokal untuk pengantaran gesit."),
        ("Transaksi Aman", "Sistem auth modern, notifikasi toast, dan konfirmasi pembayaran."),
        ("Rating & Review", "Badge, counter animasi, dan testimoni transparan."),
        ("Promo & Gamifikasi", "Badge pencapaian, promo harian, dan campaign komunitas."),
        ("Konten Edukatif", "Artikel, greeting personal, dan onboarding UMKM step-by-step."),
    ]
    card_width = Inches(4)
    card_height = Inches(2.1)
    start_left = Inches(0.8)
    start_top = Inches(2.2)
    for idx, (title, desc) in enumerate(feature_cards):
        row = idx // 3
        col = idx % 3
        left = start_left + col * (card_width + Inches(0.4))
        top = start_top + row * (card_height + Inches(0.4))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_width, card_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["cream"] if row == 0 else colors["secondary"]
        shape.line.color.rgb = colors["primary"]
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.name = "Poppins"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors["deep_green"]
        desc_p = tf.add_paragraph()
        desc_p.text = desc
        desc_p.font.name = "Nunito"
        desc_p.font.size = Pt(14)
        desc_p.font.color.rgb = colors["text_dark"]

    # Slide 6 - Ecosystem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, colors["cream"])
    add_title(slide, "Alur Ekosistem Pengguna")
    actors = [
        ("UMKM", "Kelola katalog, promo, stok, dan pesanan via dashboard."),
        ("Pengguna", "Jelajah kategori, simpan favorit, transaksi aman."),
        ("Driver", "Terima dan eksekusi tugas antar lokal real-time."),
        ("Admin", "Pantau performa, verifikasi UMKM, kelola konten."),
    ]
    card_w = Inches(3.1)
    card_h = Inches(2.5)
    for idx, (title, desc) in enumerate(actors):
        left = Inches(0.8) + idx * (card_w + Inches(0.3))
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.1), card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["white"]
        shape.line.color.rgb = colors["accent_green"]
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.name = "Poppins"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = colors["accent_dark_green"]
        desc_p = tf.add_paragraph()
        desc_p.text = desc
        desc_p.font.name = "Nunito"
        desc_p.font.size = Pt(16)
        desc_p.font.color.rgb = colors["text_dark"]

    # Slide 7 - Traction
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, colors["white"])
    add_title(slide, "Traksi & Bukti Sosial")
    stats = [
        ("500+", "UMKM terdaftar"),
        ("10K+", "Pengguna aktif per bulan"),
        ("50+", "Driver lokal siap jalan"),
        ("4.8/5", "Rating pengalaman belanja"),
    ]
    box_w = Inches(3.1)
    box_h = Inches(2.3)
    for idx, (stat, detail) in enumerate(stats):
        row = idx // 2
        col = idx % 2
        left = Inches(1) + col * (box_w + Inches(0.6))
        top = Inches(2.2) + row * (box_h + Inches(0.4))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["cream"]
        shape.line.color.rgb = colors["secondary"]
        tf = shape.text_frame
        tf.text = stat
        p = tf.paragraphs[0]
        p.font.name = "Poppins"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        detail_p = tf.add_paragraph()
        detail_p.text = detail
        detail_p.font.name = "Nunito"
        detail_p.font.size = Pt(16)
        detail_p.font.color.rgb = colors["deep_green"]

    # Slide 8 - Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, colors["cream"])
    add_title(slide, "Roadmap Singkat 2025")
    add_list(
        slide,
        [
            "Q1: Kampanye onboarding UMKM + konten edukasi video.",
            "Q2: Loyalty & referral points untuk pengguna setia.",
            "Q3: Ekspansi kategori agro & wisata kolaborasi pemda.",
            "Q4: Insight AI untuk rekomendasi produk & prediksi demand.",
        ],
    )

    # Slide 9 - Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    apply_gradient(
        background.fill,
        angle=0,
        stops=[(0.0, "primary"), (0.5, "accent_green"), (1.0, "accent_dark_green")],
    )
    background.line.fill.background()

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.2), prs.slide_width - Inches(2), Inches(5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = colors["white"]
    box.line.color.rgb = colors["secondary"]
    box.line.width = Pt(3)

    text_box = slide.shapes.add_textbox(Inches(1.4), Inches(1.6), prs.slide_width - Inches(2.8), Inches(4.2))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Mari berkolaborasi menghadirkan kebanggaan lokal"
    p.font.name = "Poppins"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors["deep_green"]

    p2 = tf.add_paragraph()
    p2.text = (
        "Kesiapan: desain produk, konten landing page, alur onboarding, dan backend siap scale. "
        "Dukungan Anda mempercepat kampanye, promosi, dan layanan baru."
    )
    p2.font.name = "Nunito"
    p2.font.size = Pt(20)
    p2.font.color.rgb = colors["text_dark"]
    p2.space_before = Pt(12)

    p3 = tf.add_paragraph()
    p3.text = "Kontak: business@aslibogor.id | +62 812-3456-7890"
    p3.font.name = "Nunito"
    p3.font.size = Pt(18)
    p3.font.color.rgb = colors["primary"]
    p3.space_before = Pt(20)

    return prs


def main() -> None:
    prs = create_presentation()
    output_dir = Path("docs")
    output_dir.mkdir(exist_ok=True)
    output_path = output_dir / "Asli_Bogor_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Saved pitch deck to {output_path}")


if __name__ == "__main__":
    main()

